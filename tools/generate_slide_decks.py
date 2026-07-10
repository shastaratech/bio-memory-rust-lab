#!/usr/bin/env python3
"""Generate Bio Memory Rust Lab slide decks from curated chapter outlines."""

from __future__ import annotations

from pathlib import Path
from textwrap import wrap

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "slides" / "book-chapters"

WIDE_W = Inches(13.333)
WIDE_H = Inches(7.5)

NAVY = RGBColor(23, 32, 51)
BLUE = RGBColor(53, 80, 112)
GREEN = RGBColor(47, 148, 109)
ORANGE = RGBColor(217, 144, 47)
LIGHT = RGBColor(248, 250, 252)
MUTED = RGBColor(81, 96, 112)
WHITE = RGBColor(255, 255, 255)
CODE_BG = RGBColor(238, 242, 247)


def add_textbox(slide, x, y, w, h, text, size=24, color=NAVY, bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    return box


def add_wrapped_text(slide, x, y, w, h, text, size=22, color=NAVY, bold=False):
    wrapped = "\n".join(wrap(text, width=max(28, int(w * 11))))
    return add_textbox(slide, x, y, w, h, wrapped, size=size, color=color, bold=bold)


def add_header(slide, chapter, title):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = LIGHT

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), WIDE_W, Inches(0.72)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    add_textbox(slide, 0.45, 0.18, 3.0, 0.35, "ShastaraTech", size=15, color=WHITE, bold=True)
    add_textbox(slide, 3.2, 0.18, 7.8, 0.35, f"{chapter} - {title}", size=14, color=WHITE)


def add_footer(slide, footer):
    add_textbox(slide, 0.45, 7.08, 8.0, 0.25, footer, size=10, color=MUTED)


def add_title_slide(prs, chapter, title, subtitle, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = LIGHT

    left = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(4.2), WIDE_H
    )
    left.fill.solid()
    left.fill.fore_color.rgb = NAVY
    left.line.fill.background()

    add_textbox(slide, 0.5, 0.55, 3.2, 0.4, "Bio Memory Rust Lab", size=18, color=WHITE, bold=True)
    add_textbox(slide, 0.5, 5.85, 3.1, 0.8, chapter, size=44, color=WHITE, bold=True)

    add_textbox(slide, 4.7, 1.35, 7.7, 1.35, title, size=38, color=NAVY, bold=True)
    add_wrapped_text(slide, 4.75, 3.0, 7.3, 1.05, subtitle, size=22, color=MUTED)
    add_textbox(slide, 4.75, 6.72, 6.0, 0.25, footer, size=11, color=MUTED)
    return slide


def add_bullets(slide, items, x=0.8, y=1.55, w=11.8, h=5.2, size=24, color=NAVY):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def add_two_column(slide, left_title, left_items, right_title, right_items):
    add_textbox(slide, 0.75, 1.15, 5.5, 0.4, left_title, size=24, color=BLUE, bold=True)
    add_textbox(slide, 6.85, 1.15, 5.5, 0.4, right_title, size=24, color=GREEN, bold=True)
    add_bullets(slide, left_items, x=0.78, y=1.72, w=5.45, h=4.7, size=19)
    add_bullets(slide, right_items, x=6.88, y=1.72, w=5.45, h=4.7, size=19)


def add_code(slide, code, x=0.75, y=1.35, w=11.9, h=4.9, size=18):
    rect = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = CODE_BG
    rect.line.color.rgb = RGBColor(203, 213, 225)
    box = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.22), Inches(w - 0.5), Inches(h - 0.35))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = code
    p.font.name = "Menlo"
    p.font.size = Pt(size)
    p.font.color.rgb = NAVY


def add_cards(slide, cards):
    x_positions = [0.7, 4.55, 8.4]
    y = 1.55
    for idx, (title, body, color) in enumerate(cards):
        x = x_positions[idx % 3]
        rect = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.25), Inches(3.75)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = WHITE
        rect.line.color.rgb = color
        rect.line.width = Pt(2)
        add_textbox(slide, x + 0.25, y + 0.25, 2.8, 0.4, title, size=22, color=color, bold=True)
        add_wrapped_text(slide, x + 0.25, y + 0.9, 2.75, 2.25, body, size=17, color=NAVY)


def normal_slide(prs, chapter, chapter_title, slide_title, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, chapter, chapter_title)
    add_textbox(slide, 0.7, 0.92, 11.8, 0.5, slide_title, size=30, color=NAVY, bold=True)
    add_footer(slide, footer)
    return slide


def make_deck(deck):
    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H
    chapter = deck["chapter"]
    title = deck["title"]
    footer = f"{deck['slug']} | exercises/rust-molecule-model"
    add_title_slide(prs, chapter, title, deck["subtitle"], footer)

    for spec in deck["slides"]:
        slide = normal_slide(prs, chapter, title, spec["title"], footer)
        kind = spec.get("kind", "bullets")
        if kind == "bullets":
            add_bullets(slide, spec["items"])
        elif kind == "two":
            add_two_column(slide, spec["left_title"], spec["left"], spec["right_title"], spec["right"])
        elif kind == "code":
            add_code(slide, spec["code"], size=spec.get("size", 18))
        elif kind == "cards":
            add_cards(slide, spec["cards"])
        else:
            raise ValueError(f"unknown slide kind: {kind}")

    out_dir = OUT / deck["slug"]
    out_dir.mkdir(parents=True, exist_ok=True)
    out_file = out_dir / f"{deck['slug']}.pptx"
    prs.save(out_file)
    return out_file


DECKS = [
    {
        "slug": "00-setup",
        "chapter": "Setup",
        "title": "Compile, Test, Run",
        "subtitle": "Get the repo running before learning chemistry or Rust concepts.",
        "slides": [
            {"title": "Student Goal", "items": [
                "Install or verify Rust and Cargo.",
                "Clone the repo and find the Rust exercise.",
                "Run tests before changing code.",
                "Run the molecule CLI and read its output.",
            ]},
            {"title": "Files To Open", "items": [
                "book/00-setup.md",
                "exercises/rust-molecule-model/Cargo.toml",
                "exercises/rust-molecule-model/src/main.rs",
                "exercises/rust-molecule-model/src/molecule.rs",
            ]},
            {"title": "Compile And Test", "kind": "code", "code": "cd exercises/rust-molecule-model\ncargo fmt --check\ncargo test"},
            {"title": "Run The CLI", "kind": "code", "code": "cargo run -- list\ncargo run -- water summary\ncargo run -- methane summary\ncargo run -- ethanol summary"},
            {"title": "Expected First Output", "items": [
                "list prints: water, methane, ethanol",
                "water formula is H2O.",
                "methane formula is CH4.",
                "ethanol formula is C2H6O.",
                "All validation checks should print true.",
            ]},
            {"title": "Checkpoint", "items": [
                "What command runs tests?",
                "What command lists toy molecules?",
                "Which file contains struct Molecule?",
            ]},
        ],
    },
    {
        "slug": "00-chemistry-primer",
        "chapter": "Chapter 00",
        "title": "Atoms, Elements, Bonds, Molecules, And Formulas",
        "subtitle": "The chemistry vocabulary students need before Rust data structures.",
        "slides": [
            {"title": "Four Core Words", "kind": "cards", "cards": [
                ("Atom", "A tiny unit of matter. In our model, it becomes one node.", BLUE),
                ("Element", "A kind of atom, such as H, C, N, or O.", GREEN),
                ("Bond", "A connection between atoms. In our model, it becomes an edge.", ORANGE),
            ]},
            {"title": "What Is A Molecule?", "items": [
                "A molecule is a group of atoms connected by bonds.",
                "Water is one oxygen atom connected to two hydrogen atoms.",
                "Methane is one carbon atom connected to four hydrogen atoms.",
                "Ethanol is a larger connected atom-bond graph.",
            ]},
            {"title": "How To Read H2O", "kind": "two", "left_title": "Formula", "left": [
                "H2 means two hydrogen atoms.",
                "O means one oxygen atom.",
                "No number means one.",
            ], "right_title": "CLI", "right": [
                "cargo run -- water summary",
                "formula: H2O",
                "atoms: 3",
                "bonds: 2",
            ]},
            {"title": "How To Read CH4", "kind": "two", "left_title": "Formula", "left": [
                "C means one carbon atom.",
                "H4 means four hydrogen atoms.",
                "Total atoms: five.",
            ], "right_title": "CLI", "right": [
                "cargo run -- methane summary",
                "formula: CH4",
                "atoms: 5",
                "bonds: 4",
            ]},
            {"title": "How To Read C2H6O", "kind": "two", "left_title": "Formula", "left": [
                "C2 means two carbon atoms.",
                "H6 means six hydrogen atoms.",
                "O means one oxygen atom.",
            ], "right_title": "CLI", "right": [
                "cargo run -- ethanol summary",
                "formula: C2H6O",
                "atoms: 9",
                "bonds: 8",
            ]},
            {"title": "Formula vs Graph", "items": [
                "A formula preserves element counts.",
                "A formula leaves out exact connectivity.",
                "A graph stores atoms plus bonds.",
                "Rust needs Atom, Bond, and Molecule because formula text is not enough.",
            ]},
            {"title": "Run It", "kind": "code", "code": "cd exercises/rust-molecule-model\ncargo run -- water summary\ncargo run -- methane summary\ncargo run -- ethanol summary"},
            {"title": "Checkpoint", "items": [
                "What does the 2 mean in H2O?",
                "What does CH4 mean?",
                "What information does C2H6O leave out?",
                "Why does Rust need a Bond type?",
            ]},
        ],
    },
    {
        "slug": "01-atoms-memory-types",
        "chapter": "Chapter 01",
        "title": "Atoms, Memory, And Types",
        "subtitle": "Use Rust types to make atom facts explicit and reliable.",
        "slides": [
            {"title": "Big Idea", "items": [
                "A type is a rule for what kind of value something is.",
                "An atom has typed facts: element, charge, aromatic flag.",
                "Rust catches many representation mistakes early.",
            ]},
            {"title": "Element As Enum", "kind": "code", "code": "pub enum Element {\n    H,\n    C,\n    N,\n    O,\n    F,\n    Cl,\n    Br,\n}"},
            {"title": "Atom As Struct", "kind": "code", "code": "pub struct Atom {\n    pub element: Element,\n    pub formal_charge: i8,\n    pub aromatic: bool,\n}"},
            {"title": "Why Not Just Strings?", "kind": "two", "left_title": "String Risk", "left": [
                '"cl" and "CL" can both appear.',
                "Misspellings are runtime data problems.",
                "Text is useful at file boundaries.",
            ], "right_title": "Enum Benefit", "right": [
                "Element::Cl has one spelling.",
                "Compiler checks allowed variants.",
                "Clean internal model.",
            ]},
            {"title": "Compile And Inspect", "kind": "code", "code": "cd exercises/rust-molecule-model\ncargo test\ncargo run -- water atoms\ncargo run -- water bonds"},
            {"title": "Checkpoint", "items": [
                "Why is Element an enum?",
                "Why is formal_charge signed?",
                "What does Atom::neutral(Element::O) create?",
            ]},
        ],
    },
    {
        "slug": "02-molecules-as-graphs",
        "chapter": "Chapter 02",
        "title": "Molecules As Graphs",
        "subtitle": "Atoms become nodes, bonds become edges, and algorithms become possible.",
        "slides": [
            {"title": "Graph Vocabulary", "items": [
                "Node: an atom.",
                "Edge: a bond.",
                "Neighbor: directly bonded atom.",
                "Path: route through bonds.",
                "Component: separate fragment.",
            ]},
            {"title": "Storage vs View", "kind": "two", "left_title": "Stored", "left": [
                "Vec<Atom>",
                "Vec<Bond>",
                "Simple and readable.",
            ], "right_title": "Computed View", "right": [
                "Vec<Vec<usize>> adjacency list",
                "Fast repeated neighbor lookup.",
                "Good for BFS.",
            ]},
            {"title": "Neighbor Query", "kind": "code", "code": "cargo run -- ethanol atoms\ncargo run -- ethanol neighbors 1\n\n# expected neighbors:\n[0, 2, 6, 7]"},
            {"title": "Shortest Path", "kind": "code", "code": "cargo run -- ethanol path 3 8\n\n# expected path:\n[3, 0, 1, 2, 8]"},
            {"title": "BFS Intuition", "items": [
                "Start at one atom.",
                "Visit all unvisited neighbors.",
                "Remember where each atom came from.",
                "Stop when the goal is found.",
                "Reconstruct the path.",
            ]},
            {"title": "Checkpoint", "items": [
                "What is a node?",
                "What is an edge?",
                "Why can a path be missing?",
                "Why use VecDeque for BFS?",
            ]},
        ],
    },
    {
        "slug": "03-ownership-borrowing-design",
        "chapter": "Chapter 03",
        "title": "Ownership And Borrowing Through Molecule Design",
        "subtitle": "Many readers or one editor: Rust's sharing rule as a molecule workflow.",
        "slides": [
            {"title": "Big Idea", "items": [
                "One owner holds the molecule data.",
                "Many readers can inspect it.",
                "One editor can change it.",
                "Validation keeps broken molecules from moving forward.",
            ]},
            {"title": "Borrowed Methods", "kind": "code", "code": "pub fn atoms(&self) -> &[Atom]\npub fn bonds(&self) -> &[Bond]\npub fn validate_bond_indices(&self) -> bool\npub fn formula(&self) -> String"},
            {"title": "Run Validation", "kind": "code", "code": "cargo run -- water validate\ncargo run -- ethanol validate\n\n# expected:\ntrue\ntrue"},
            {"title": "Why Borrow?", "items": [
                "Summary can read without owning the molecule.",
                "Formula can read without copying all atoms.",
                "Neighbor queries can reuse the same molecule.",
                "Scoring and validation can be composed safely.",
            ]},
            {"title": "Student Edit", "items": [
                "Open src/main.rs.",
                "Find print_summary.",
                "Add a line for connected_components().",
                "Run cargo run -- ethanol summary.",
                "Run cargo test.",
            ]},
            {"title": "Checkpoint", "items": [
                "What does &self mean?",
                "Why does print_summary use &Molecule?",
                "What does validation check?",
            ]},
        ],
    },
    {
        "slug": "04-quantum-state-simulation",
        "chapter": "Chapter 04",
        "title": "Quantum State And Molecular Simulation",
        "subtitle": "A graph is useful, but a molecule is more than a graph.",
        "slides": [
            {"title": "What The Graph Can Answer", "items": [
                "Formula",
                "Atom count",
                "Bond count",
                "Neighbors",
                "Paths",
                "Connected components",
            ]},
            {"title": "What The Graph Leaves Out", "items": [
                "Exact 3D shape.",
                "Electron density.",
                "Real reaction behavior.",
                "Quantum energy.",
                "Docking score.",
            ]},
            {"title": "Run The Graph Model", "kind": "code", "code": "cargo run -- ethanol summary\ncargo run -- ethanol path 3 8"},
            {"title": "Quantum-Classical Loop", "items": [
                "Choose molecular problem.",
                "Build simplified Hamiltonian.",
                "Prepare quantum state.",
                "Measure energy terms.",
                "Classical optimizer updates parameters.",
                "Repeat until energy stops improving.",
            ]},
            {"title": "Guardrail", "kind": "two", "left_title": "Avoid", "left": [
                "A qubit is a bigger bit.",
                "Quantum stores all answers.",
                "Quantum solves all chemistry.",
            ], "right_title": "Say Instead", "right": [
                "A qubit is a quantum state.",
                "Algorithms shape amplitudes.",
                "Quantum may help some hard simulations.",
            ]},
            {"title": "Checkpoint", "items": [
                "Why is a graph not a complete molecule?",
                "Which part of VQE is classical?",
                "Which part prepares and measures quantum state?",
            ]},
        ],
    },
    {
        "slug": "05-data-structures",
        "chapter": "Chapter 05",
        "title": "Rust Data Structures For Molecules",
        "subtitle": "Choose structures based on the question you want to answer.",
        "slides": [
            {"title": "Representation Choices", "kind": "two", "left_title": "Question", "left": [
                "What elements are allowed?",
                "What facts belong to an atom?",
                "How many atoms exist?",
                "Which atoms are connected?",
            ], "right_title": "Structure", "right": [
                "enum Element",
                "struct Atom",
                "Vec<Atom>",
                "Vec<Bond> or adjacency list",
            ]},
            {"title": "Molecule Owns Collections", "kind": "code", "code": "pub struct Molecule {\n    atoms: Vec<Atom>,\n    bonds: Vec<Bond>,\n}"},
            {"title": "Formula Counts", "kind": "code", "code": "pub fn formula_counts(&self) -> Vec<(Element, usize)>"},
            {"title": "Run Formula Commands", "kind": "code", "code": "cargo run -- water formula\ncargo run -- methane formula\ncargo run -- ethanol formula"},
            {"title": "Try It: Add Ammonia", "items": [
                "Add pub fn ammonia() -> Molecule.",
                "Use one N and three H atoms.",
                "Add three N-H bonds.",
                "Update molecule_by_name.",
                "Update molecule_names.",
                "Run cargo test.",
            ]},
            {"title": "Checkpoint", "items": [
                "Why use Vec<Atom>?",
                "What does Vec<(Element, usize)> represent?",
                "What is storage representation vs graph view?",
            ]},
        ],
    },
    {
        "slug": "06-traits-functions",
        "chapter": "Chapter 06",
        "title": "Traits And Functions For Molecular Behavior",
        "subtitle": "Types describe data; traits describe shared capabilities.",
        "slides": [
            {"title": "Three Traits", "kind": "cards", "cards": [
                ("Describe", "Atom, Bond, and Molecule can explain themselves.", BLUE),
                ("ChemicalFormula", "Molecule can produce formula counts and formula text.", GREEN),
                ("MolecularGraph", "Molecule can answer graph questions.", ORANGE),
            ]},
            {"title": "Describe Trait", "kind": "code", "code": "pub trait Describe {\n    fn describe(&self) -> String;\n}"},
            {"title": "Formula Trait", "kind": "code", "code": "pub trait ChemicalFormula {\n    fn formula_counts(&self) -> Vec<(Element, usize)>;\n\n    fn formula(&self) -> String {\n        ...\n    }\n}"},
            {"title": "Graph Trait", "kind": "code", "code": "pub trait MolecularGraph {\n    fn atom_count(&self) -> usize;\n    fn bond_count(&self) -> usize;\n    fn neighbors(&self, atom_id: usize) -> Vec<usize>;\n    fn shortest_path(&self, start: usize, goal: usize) -> Option<Vec<usize>>;\n}"},
            {"title": "Run It", "kind": "code", "code": "cargo run -- water summary\ncargo run -- ethanol neighbors 1\ncargo run -- ethanol path 3 8"},
            {"title": "Checkpoint", "items": [
                "What is a trait?",
                "Which trait has a default method?",
                "Why accept &impl MolecularGraph?",
            ]},
        ],
    },
    {
        "slug": "07-molecule-explorer-project",
        "chapter": "Chapter 07",
        "title": "Molecule Explorer CLI Project",
        "subtitle": "Turn the molecule model into a tool students can run and extend.",
        "slides": [
            {"title": "Project Goal", "items": [
                "Inspect toy molecules from the command line.",
                "Connect commands to Rust functions.",
                "Practice library plus CLI structure.",
                "Add a new command or molecule.",
            ]},
            {"title": "File Structure", "kind": "code", "code": "src/lib.rs       # reusable library\nsrc/main.rs      # command-line interface\nsrc/molecule.rs  # data model, traits, algorithms"},
            {"title": "Run All Main Commands", "kind": "code", "code": "cargo run -- list\ncargo run -- water summary\ncargo run -- ethanol atoms\ncargo run -- ethanol bonds\ncargo run -- ethanol neighbors 1\ncargo run -- ethanol path 3 8"},
            {"title": "Command Flow", "items": [
                "Read command-line args.",
                "Look up molecule by name.",
                "Match the requested command.",
                "Call Molecule methods.",
                "Print student-readable output.",
            ]},
            {"title": "Project Task: Add degree", "kind": "code", "code": "cargo run -- ethanol degree 1\n\n# expected:\n4\n\n# hint:\nmolecule.neighbors(atom_id).len()"},
            {"title": "University Extension", "items": [
                "Introduce enum Command.",
                "Write parse_command(args) -> Result<Command, String>.",
                "Unit test command parsing.",
                "Keep CLI messages friendly.",
            ]},
            {"title": "Checkpoint", "items": [
                "What is the difference between lib.rs and main.rs?",
                "Which file stores algorithms?",
                "How does the CLI choose a molecule?",
            ]},
        ],
    },
]


def write_index(outputs):
    lines = [
        "# Slide Decks",
        "",
        "Generated PowerPoint decks for the Bio Memory Rust Lab self-study book.",
        "",
        "## Decks",
        "",
        "| Chapter | PPTX | Source chapter |",
        "| --- | --- | --- |",
    ]
    source_map = {
        "00-setup": "book/00-setup.md",
        "00-chemistry-primer": "book/00-chemistry-primer.md",
        "01-atoms-memory-types": "book/01-atoms-memory-types.md",
        "02-molecules-as-graphs": "book/02-molecules-as-graphs.md",
        "03-ownership-borrowing-design": "book/03-ownership-borrowing-design.md",
        "04-quantum-state-simulation": "book/04-quantum-state-simulation.md",
        "05-data-structures": "book/05-data-structures.md",
        "06-traits-functions": "book/06-traits-functions.md",
        "07-molecule-explorer-project": "book/07-molecule-explorer-project.md",
    }
    for deck, output in outputs:
        rel = output.relative_to(ROOT / "slides")
        src = source_map[deck["slug"]]
        lines.append(f"| {deck['chapter']}: {deck['title']} | [{rel.name}]({rel}) | [{src}](../{src}) |")
    lines += [
        "",
        "## Regenerate",
        "",
        "```bash",
        "python3 tools/generate_slide_decks.py",
        "```",
        "",
        "The PPTX files are generated from curated outlines in `tools/generate_slide_decks.py`.",
    ]
    (ROOT / "slides" / "README.md").write_text("\n".join(lines) + "\n")


def main():
    outputs = []
    for deck in DECKS:
        outputs.append((deck, make_deck(deck)))
    write_index(outputs)
    for _, path in outputs:
        print(path.relative_to(ROOT))


if __name__ == "__main__":
    main()
